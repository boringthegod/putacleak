import argparse
import os
import re
import urllib.parse
import time
from collections import defaultdict
import requests
import fitz
import pandas as pd
from openpyxl import load_workbook
from odf.opendocument import load
from odf.text import P
from docx import Document
import csv 
import docx2txt
import subprocess
import warnings
from pptx import Presentation
import chardet
import ollama
from datetime import datetime
import re

warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")

proxies = {
    # "http": "socks5://your_proxy_here",
    # "https": "socks5://your_proxy_here"
}

headers = {
    'Host': 'www.google.fr',
    'Sec-Ch-Ua': '"Not?A_Brand";v="99", "Chromium";v="130"',
    'Sec-Ch-Ua-Mobile': '?0',
    'Sec-Ch-Ua-Platform': '"Windows"',
    'Sec-Ch-Prefers-Color-Scheme': 'light',
    'Accept-Language': 'fr-FR,fr;q=0.9',
    'Upgrade-Insecure-Requests': '1',
    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/130.0.6723.59 Safari/537.36',
    'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8',
    'Sec-Fetch-Site': 'same-origin',
    'Sec-Fetch-Mode': 'navigate',
    'Sec-Fetch-User': '?1',
    'Sec-Fetch-Dest': 'document',
    'Referer': 'https://www.google.fr/',
    'Priority': 'u=0, i',
}

default_keywords = '"mot de passe" OR password OR passwort OR secrets OR secret OR confidentiel OR confidential OR token OR passwd OR aws_access OR credentials OR login OR "mot de passes" OR passwords OR "internal use only" OR administrator'
default_filetypes = 'filetype:doc OR filetype:docx OR filetype:pdf OR filetype:xls OR filetype:xlsx OR filetype:txt OR filetype:ppt OR filetype:odt OR filetype:csv'

additional_dorks = [
    ('api_secret OR access_key OR token OR passwd OR password OR passwords OR passwort OR secrets OR secret OR aws_access OR login OR pass OR credentials OR admin',
     'filetype:ps1 OR filetype:sql OR filetype:sh OR filetype:bash OR filetype:xml OR filetype:conf OR filetype:config OR filetype:bak'),

    ('api_secret OR access_key OR token OR passwd OR password OR passwords OR passwort OR secrets OR secret OR aws_access OR login OR pass OR credentials OR admin',
     'filetype:cnf OR filetype:log OR filetype:htaccess OR filetype:env OR filetype:ini OR filetype:pwd OR filetype:json OR filetype:logs OR filetype:old OR filetype:yml OR filetype:yaml OR filetype:cgi')
]

successful_downloads = defaultdict(int)
failed_downloads = []

def extract_and_download_urls(response_text, download_path, verbose):
    urls = re.findall(r'url=(.*?)&amp;', response_text)
    for url in urls:
        decoded_url = urllib.parse.unquote(url)
        if not decoded_url.startswith('/search?q=') and \
           not decoded_url.startswith("https://support.google.com") and \
           not decoded_url.startswith("https://maps.google.fr") and \
           "enablejs" not in decoded_url:
            if verbose:
                print(f"Download of : {decoded_url}")
            download_file(decoded_url, download_path, verbose)

def download_file(url, download_path, verbose, retries=2):
    for attempt in range(1, retries + 1):
        try:
            if attempt == 1:
                response = requests.get(url, verify=False)
                if verbose:
                    print(f"Download attempt without proxy : {url}")
            else:
                response = requests.get(url, proxies=proxies, verify=False)
                if verbose:
                    print(f"Download attempt with proxy (attempt {attempt}/{retries}): {url}")
            
            response.raise_for_status()
            filename = os.path.join(download_path, os.path.basename(urllib.parse.urlparse(url).path))
            file_extension = os.path.splitext(filename)[1].lower().replace(".", "")
            
            with open(filename, 'wb') as file:
                file.write(response.content)
            
            successful_downloads[file_extension] += 1
            if verbose:
                print(f"Downloaded file : {filename}")
            return

        except Exception as e:
            if attempt == retries:
                failed_downloads.append((url, str(e)))
            elif verbose:
                print(f"Error while downloading {url} : {e}")
                print("New attempt...")

def fetch_page(params, verbose, retries=3):
    attempt = 0
    while attempt < retries:
        try:
            response = requests.get(
                'https://www.google.fr/search',
                params=params,
                headers=headers,
                proxies=proxies,
                timeout=15
            )
            response.raise_for_status()
            if "enablejs" in response.text or "Cliquez ici si" in response.text:
                if verbose:
                    print("Robot protection detected, try again.")
                attempt += 1
                time.sleep(2)
            else:
                return response
        except Exception as e:
            if verbose:
                print(f"Error during attempt {attempt + 1}/{retries}: {e}")
            attempt += 1
            time.sleep(2)
    return None

def fetch_all_results(domaine, keywords, filetypes, max_pages, use_additional_dorks, verbose):
    download_path = domaine
    os.makedirs(download_path, exist_ok=True)

    custom_dork = (keywords, filetypes)
    dorks_to_use = [custom_dork]

    if use_additional_dorks:
        dorks_to_use += additional_dorks

    for keywords, filetypes in dorks_to_use:
        query = f'site:{domaine} ({keywords}) ({filetypes})'
        print(f"\nSearch in progress for keywords : {keywords}")
        print(f"File types : {filetypes}\n")

        for page_num in range(max_pages):
            start_param = page_num * 10
            params = {'q': query + f'&start={start_param}'}

            response = fetch_page(params, verbose)
            if response:
                extract_and_download_urls(response.text, download_path, verbose)
                if not re.search(rf'></span>{page_num+2}</a></td><td', response.text):
                    if verbose:
                        print(f"No page {page_num + 2}, search stopped for this dork.")
                    break

def extract_text_from_pdf(pdf_path, output_txt_path):
    with fitz.open(pdf_path) as pdf:
        with open(output_txt_path, "w", encoding="utf-8") as output_file:
            for page_num in range(pdf.page_count):
                page = pdf[page_num]
                text = page.get_text("text")
                output_file.write(text)
                output_file.write("\n\n")
    print(f"Text extracted and saved in {output_txt_path}")

def excel_to_text(excel_path, output_txt_path):
    xls = pd.ExcelFile(excel_path)
    with open(output_txt_path, "w", encoding="utf-8") as output_file:
        for sheet_name in xls.sheet_names:
            output_file.write(f"\n--- Sheet: {sheet_name} ---\n")
            sheet = pd.read_excel(xls, sheet_name=sheet_name, header=None)
            for _, row in sheet.iterrows():
                row_text = "\t".join([str(cell) if not pd.isna(cell) else "" for cell in row])
                output_file.write(row_text + "\n")
    print(f"Contents of Excel file saved in {output_txt_path}")

def extract_text_from_doc(doc_path, output_txt_path):
    try:
        text = docx2txt.process(doc_path)
        with open(output_txt_path, "w", encoding="utf-8") as output_file:
            output_file.write(text)
        print(f"Text extracted and saved in {output_txt_path}")
    except Exception as e:
        print(f"Error extracting {doc_path} with docx2txt:: {e}")
        try:
            text = subprocess.check_output(['antiword', doc_path], universal_newlines=True)
            with open(output_txt_path, "w", encoding="utf-8") as output_file:
                output_file.write(text)
            print(f"Text extracted with antiword and saved in {output_txt_path}")
        except FileNotFoundError:
            print("antiword is not installed. Install it to manage older .doc files.")

def convert_docm_to_docx(docm_path, output_docx_path):
    try:
        subprocess.run(["libreoffice", "--headless", "--convert-to", "docx", docm_path, "--outdir", os.path.dirname(output_docx_path)], check=True)
        print(f"File {docm_path} converted to {output_docx_path}")
        return output_docx_path
    except Exception as e:
        print(f"Conversion error for {docm_path}: {e}")
        return None

def extract_text_from_pptx(pptx_path, output_txt_path):
    presentation = Presentation(pptx_path)
    with open(output_txt_path, "w", encoding="utf-8") as output_file:
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    output_file.write(shape.text + "\n")
    print(f"Text extracted and saved in {output_txt_path}")

def extract_text_from_xlsx(xlsx_path, output_txt_path):
    workbook = load_workbook(xlsx_path, data_only=True)
    with open(output_txt_path, "w", encoding="utf-8") as output_file:
        for sheet in workbook.sheetnames:
            output_file.write(f"\n--- Sheet: {sheet} ---\n")
            worksheet = workbook[sheet]
            for row in worksheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                output_file.write(row_text + "\n")
    print(f"Text extracted and saved in {output_txt_path}")


def extract_text_from_docx(docx_path, output_txt_path):
    doc = Document(docx_path)
    with open(output_txt_path, "w", encoding="utf-8") as output_file:
        for paragraph in doc.paragraphs:
            output_file.write(paragraph.text + "\n")
    print(f"Text extracted and saved in {output_txt_path}")

def odt_to_text(odt_path, output_txt_path):
    doc = load(odt_path)

    with open(output_txt_path, "w", encoding="utf-8") as output_file:
        for paragraph in doc.getElementsByType(P):
            text = ''.join(node.data for node in paragraph.childNodes if node.nodeType == 3)
            output_file.write(text + "\n")

    print(f"Text extracted and saved in {output_txt_path}")

def csv_to_text(csv_path, output_txt_path, delimiter="\t"):
    with open(csv_path, "rb") as f:
        result = chardet.detect(f.read())
        encoding = result['encoding'] if result['confidence'] > 0.7 else 'utf-8'

    with open(csv_path, "r", encoding=encoding) as csv_file:
        csv_reader = csv.reader(csv_file, delimiter=delimiter)
        with open(output_txt_path, "w", encoding="utf-8") as txt_file:
            for row in csv_reader:
                row_text = delimiter.join(row)
                txt_file.write(row_text + "\n")

    print(f"Text extracted and saved in {output_txt_path}")

def process_file(file_path, output_folder):
    file_name, file_ext = os.path.splitext(file_path)
    output_txt_path = os.path.join(output_folder, os.path.basename(file_name) + ".txt")

    extensions_to_convert = {
        'pdf': 'pdf', 'xls': 'xls', 'xlsx': 'xlsx', 'odt': 'odt', 'docx': 'docx', 'csv': 'csv', 'pptx': 'pptx', 'doc': 'doc'
    }

    readable_extensions = [
        'txt', 'ps1', 'sql', 'sh', 'bash', 'xml', 'conf', 'config', 'bak', 'cnf',
        'log', 'htaccess', 'env', 'ini', 'ppt', 'pwd', 'json', 'logs', 'old', 'rdp', 'yml', 'yaml', 'cgi'
    ]

    if file_ext[1:].lower() in readable_extensions:
        with open(file_path, "r", encoding="utf-8") as f_in, open(output_txt_path, "w", encoding="utf-8") as f_out:
            f_out.write(f_in.read())
        print(f"File {file_path} copied without conversion to {output_txt_path}")

    elif file_ext[1:].lower() == "pdf":
        extract_text_from_pdf(file_path, output_txt_path)
    elif file_ext[1:].lower() == "xls":
        excel_to_text(file_path, output_txt_path)
    elif file_ext[1:].lower() == "xlsx":
        extract_text_from_xlsx(file_path, output_txt_path)
    elif file_ext[1:].lower() == "odt":
        odt_to_text(file_path, output_txt_path)
    elif file_ext[1:].lower() == "docx":
        extract_text_from_docx(file_path, output_txt_path)
    elif file_ext[1:].lower() == "doc":
        extract_text_from_doc(file_path, output_txt_path)
    elif file_ext[1:].lower() == "pptx":
        extract_text_from_pptx(file_path, output_txt_path)
    elif file_ext[1:].lower() == "csv":
        csv_to_text(file_path, output_txt_path)
    elif file_ext[1:].lower() == "docm":
        docx_path = convert_docm_to_docx(file_path, file_name + ".docx")
        if docx_path:
            extract_text_from_docx(docx_path, output_txt_path)
    else:
        print(f"File type {file_ext} not supported for {file_path}")

def process_directory(directory, output_folder):
    for root, _, files in os.walk(directory):
        for file in files:
            file_path = os.path.join(root, file)
            process_file(file_path, output_folder)

def extract_surrounding_text(file_path, keywords, context_size=150, extract_output_folder=None):
    with open(file_path, "r", encoding="utf-8") as file:
        content = file.read()

    content_lower = content.lower()
    keyword_groups = [
        ["mot de passe", "mot de passes", "password", "passwords", "passwort", "passwd", "pass"],
        ["secrets", "secret"],
        ["confidentiel", "confidential"],
        ["credentials", "credential"],
        ["aws_access", "access_key"],
        ["internal use only"],
        ["login"],
        ["administrator"],
        ["api_secret"],
        ["token"]
    ]

    surrounding_texts = []
    found_positions = set()
    for group in keyword_groups:
        for keyword in group:
            positions = [i for i in range(len(content_lower)) if content_lower.startswith(keyword, i)]
            for pos in positions:
                if not any(pos in range(fp - context_size, fp + context_size) for fp in found_positions):
                    start = max(0, pos - context_size)
                    end = min(len(content), pos + len(keyword) + context_size)
                    surrounding_texts.append(content[start:end])
                    found_positions.add(pos)

    if extract_output_folder:
        os.makedirs(extract_output_folder, exist_ok=True)
        output_file = os.path.join(extract_output_folder, f"{os.path.splitext(os.path.basename(file_path))[0]}_extract.txt")
    else:
        output_file = f"{os.path.splitext(os.path.basename(file_path))[0]}_extract.txt"

    with open(output_file, "w", encoding="utf-8") as output:
        output.write(f"Extracts from the file {file_path}:\n\n")
        for i, text in enumerate(surrounding_texts, 1):
            output.write(f"Occurrence {i}:\n{text}\n\n")

    print(f"Extract saved in {output_file}")

def analyze_occurrences_in_file(file_path, model_name="artifish/llama3.2-uncensored"):
    with open(file_path, "r", encoding="utf-8") as file:
        content = file.read()

    occurrences = content.split("Occurrence ")[1:]
    leaked_count = 0
    total_count = len(occurrences)

    for occurrence in occurrences:
        occurrence_text = occurrence.split("\n", 1)[1].strip().replace("\n", " ")

        formatted_text = f"{{{{{occurrence_text}}}}}"

        message_content = (
            f"You're an expert in pentesting. I give you a text surrounded by {{ at the beginning and }} at the end. Look carefully at all its contents and if there is a leak of password or critical information you return 'CONTENT_LEAKED' otherwise you return: 'CONTENT SAFE'"
            f"\n\nHere's the text:\n\n{formatted_text}\n\n Don't write anything else"
        )

        response = ollama.chat(model=model_name, messages=[{"role": "user", "content": message_content}])

        result = response['message']['content'].strip()

        if result == "CONTENT_LEAKED":
            leaked_count += 1

    if total_count == 0:
        leaked_percentage = 0
    else:
        leaked_percentage = (leaked_count / total_count) * 100

    if leaked_percentage > 49:
        print(f"{os.path.basename(file_path)}: CONTENT_LEAKED")
        return f"{os.path.basename(file_path)}: CONTENT_LEAKED"
    else:
        print(f"{os.path.basename(file_path)}: CONTENT SAFE")
        return None

def main():
    parser = argparse.ArgumentParser(description="putacleak - search for potentially sensitive files using google dorking and then analyze their content with LLM")
    parser.add_argument("-d", "--domain", required=True, help="Domain to search, for example: carrefour.fr")
    parser.add_argument("-ft", "--filetype", help="File types to search, separated by commas (e.g.: pdf,docx,doc)")
    parser.add_argument("-kw", "--keywords", help="Keywords to search for, separated by commas (e.g.: creds,\"mot de passe\",admin)")
    parser.add_argument("-mp", "--max-pages", type=int, default=9, help="Maximum number of pages to scrape (default 9)")
    parser.add_argument("-v", "--verbose", action="store_true", help="Activate verbose mode for more details")
    args = parser.parse_args()

    filetypes = ' OR '.join([f"filetype:{ft.strip()}" for ft in args.filetype.split(',')]) if args.filetype else default_filetypes
    keywords = ' OR '.join([kw.strip() for kw in args.keywords.split(',')]) if args.keywords else default_keywords
    use_additional_dorks = not args.keywords and not args.filetype

    fetch_all_results(args.domain, keywords, filetypes, args.max_pages, use_additional_dorks, args.verbose)

    download_directory = args.domain
    output_folder = os.path.join(args.domain, "converted_txt")
    os.makedirs(output_folder, exist_ok=True)
    process_directory(download_directory, output_folder)

    extract_output_folder = os.path.join(args.domain, "extract")

    if not os.path.exists(extract_output_folder):
        print("Aucun fichier trouvé")
        return
    
    for file_name in os.listdir(output_folder):
        extract_surrounding_text(
            os.path.join(output_folder, file_name),
            keywords=[
                "mot de passe", "mot de passes", "password", "passwords", "passwort", "secrets",
                "secret", "confidentiel", "confidential", "token", "passwd", "aws_access",
                "credentials", "login", "pass", "internal use only", "administrator", "api_secret", "access_key"
            ],
            extract_output_folder=extract_output_folder
        )

    results = []
    for file_name in os.listdir(extract_output_folder):
        if file_name.endswith("_extract.txt"):
            result = analyze_occurrences_in_file(os.path.join(extract_output_folder, file_name))
            if result:
                results.append(result)

    if results:
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        
        safe_domain = re.sub(r'[^\w\-\.]', '_', args.domain)
        
        filename = f"{timestamp}_{safe_domain}_results.txt"
        
        with open(filename, "a", encoding="utf-8") as file:
            file.write(f"\n[{timestamp}] Potential Leaks Detected for {args.domain}:\n")
            for line in results:
                file.write(f"{line}\n")

    print("\n--- Overview ---")
    print("Files downloaded successfully :")
    for ext, count in successful_downloads.items():
        print(f" - {ext.upper()} : {count} file(s)")

    print("\nFiles not downloaded :")
    for url, error in failed_downloads:
        print(f" - {url} (Error : {error})")


if __name__ == "__main__":
    main()
